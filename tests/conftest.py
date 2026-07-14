"""테스트용 이메일 픽스처 팩토리(.eml=표준 email, .msg=extract_msg OleWriter)."""

from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture
def make_eml():
    from email.message import EmailMessage

    def _make(path: Path, *, subject="", body="", sender="", to="", cc="", attachments=()):
        path.parent.mkdir(parents=True, exist_ok=True)
        msg = EmailMessage()
        if subject:
            msg["Subject"] = subject
        if sender:
            msg["From"] = sender
        if to:
            msg["To"] = to
        if cc:
            msg["Cc"] = cc
        msg.set_content(body)
        for name in attachments:
            msg.add_attachment(
                b"x", maintype="application", subtype="octet-stream", filename=name
            )
        path.write_bytes(msg.as_bytes())
        return path

    return _make


@pytest.fixture
def make_msg():
    from extract_msg import OleWriter

    def _make(path: Path, *, subject="", body="", sender=""):
        path.parent.mkdir(parents=True, exist_ok=True)
        writer = OleWriter()

        def prop(propid: int, text: str) -> None:
            if text:
                writer.addEntry(f"__substg1.0_{propid:04X}001F", text.encode("utf-16-le"))

        prop(0x001A, "IPM.Note")  # PidTagMessageClass (유효 MSG 판별에 필요)
        prop(0x0037, subject)     # PidTagSubject
        prop(0x1000, body)        # PidTagBody
        prop(0x0C1A, sender)      # PidTagSenderName
        writer.addEntry("__properties_version1.0", b"\x00" * 32)
        writer.write(str(path))
        return path

    return _make


@pytest.fixture
def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    def _make(path: Path, *, title="", body_lines=(), table=None, notes="", split_runs=None):
        path.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
        tf = tb.text_frame
        tf.text = title
        for line in body_lines:
            tf.add_paragraph().text = line
        if split_runs:  # 한 문단을 여러 런으로 분할(쪼개진 키워드 테스트용)
            para = tf.add_paragraph()
            for piece in split_runs:
                run = para.add_run()
                run.text = piece
        if table:
            rows, cols = len(table), len(table[0])
            gt = slide.shapes.add_table(rows, cols, Inches(1), Inches(4.5), Inches(5), Inches(1)).table
            for r, rowvals in enumerate(table):
                for c, val in enumerate(rowvals):
                    gt.cell(r, c).text = val
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        prs.save(str(path))
        return path

    return _make
