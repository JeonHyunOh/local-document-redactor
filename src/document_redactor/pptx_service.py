"""PowerPoint(.pptx) 슬라이드 텍스트의 키워드 검색·제거·재검증.

python-pptx 의존을 이 모듈에 격리한다. 텍스트가 여러 런(run)으로 쪼개질 수 있으므로
검색·재검증은 문단(paragraph) 단위로 합쳐 판정하고, 제거는 런 단위 후 문단에 키워드가
남으면 문단 단위로 폴백해 확실히 제거한다. 원본은 편집하지 않고 <stem>_edited.pptx로 저장한다.
"""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

from pptx import Presentation

from . import keyword_matcher, pattern_matcher
from .models import (
    EditRequest,
    EditResult,
    FileType,
    PptxMatch,
    SearchCriteria,
    SearchReport,
    VerificationResult,
)


def _iter_paragraphs(prs) -> Iterator[tuple[int, str, object]]:
    """(슬라이드 1-기반 번호, location, paragraph)를 순회한다. 도형·표·노트를 포함."""
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    yield slide_index, "본문", para
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            yield slide_index, "표", para
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                yield slide_index, "노트", para


def _count(text: str, keywords: list[str], case_sensitive: bool) -> int:
    """문단 텍스트의 키워드(부분문자열) 총 발견 횟수."""
    total = 0
    haystack = text if case_sensitive else text.casefold()
    for keyword in keyword_matcher.normalize_keywords(keywords):
        needle = keyword if case_sensitive else keyword.casefold()
        if needle:
            total += haystack.count(needle)
    return total


def _text_hits(text: str, criteria: SearchCriteria) -> int:
    """문단 텍스트의 키워드 + 패턴 매치 총 개수."""
    kw = _count(text, criteria.keywords, criteria.case_sensitive)
    pat = len(pattern_matcher.find_patterns(text, criteria.redact_account_numbers))
    return kw + pat


def _redact(text: str, criteria: SearchCriteria) -> str:
    """키워드 제거 후 패턴도 제거한다."""
    out = keyword_matcher.remove_keywords(text, criteria.keywords, criteria.case_sensitive)
    return pattern_matcher.remove_patterns(out, criteria.redact_account_numbers)


def search(path: Path, criteria: SearchCriteria) -> SearchReport:
    """슬라이드·표·노트 문단에서 키워드·패턴을 검색한다(파일 무수정)."""
    path = Path(path)
    prs = Presentation(str(path))
    keywords = keyword_matcher.normalize_keywords(criteria.keywords)
    matches: list[PptxMatch] = []
    for slide_index, location, para in _iter_paragraphs(prs):
        text = para.text
        for keyword in keywords:
            n = _count(text, [keyword], criteria.case_sensitive)
            if n:
                matches.append(
                    PptxMatch(
                        file_name=path.name,
                        slide=slide_index,
                        location=location,
                        keyword=keyword,
                        count=n,
                        context=text,
                    )
                )
        for label, _value in pattern_matcher.find_patterns(text, criteria.redact_account_numbers):
            matches.append(
                PptxMatch(
                    file_name=path.name,
                    slide=slide_index,
                    location=location,
                    keyword=f"[{label}]",
                    count=1,
                    context=text,
                )
            )
    return SearchReport(
        file_name=path.name, file_type=FileType.PPTX, criteria=criteria, pptx_matches=matches
    )


def apply_edit(
    path: Path, request: EditRequest, output_dir: Path, selected=None
) -> EditResult:
    """슬라이드 텍스트에서 키워드를 제거해 <stem>_edited.pptx로 저장한다(selected 무시)."""
    path = Path(path)
    prs = Presentation(str(path))

    removed = 0
    for _, _, para in _iter_paragraphs(prs):
        before = _text_hits(para.text, request.criteria)
        if before == 0:
            continue
        removed += before
        # 런 단위 제거(키워드 + 패턴)
        for run in para.runs:
            run.text = _redact(run.text, request.criteria)
        # 문단에 남으면(쪼개진 런) 첫 런에 정리된 문단 텍스트를 넣고 나머지 런 비우기
        if _text_hits(para.text, request.criteria) > 0 and para.runs:
            cleaned = _redact(para.text, request.criteria)
            para.runs[0].text = cleaned
            for run in para.runs[1:]:
                run.text = ""

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"{path.stem}_edited.pptx"
    prs.save(str(out_path))

    return EditResult(
        source_name=path.name,
        output_path=str(out_path),
        file_type=FileType.PPTX,
        redactions_applied=removed,
        log=[f"슬라이드 텍스트에서 키워드 {removed}건 제거 → {out_path.name}"],
    )


def verify(output_path: Path, criteria: SearchCriteria) -> VerificationResult:
    """저장된 .pptx를 문단 단위로 재검색해 키워드 잔존 여부를 확인한다."""
    remaining = search(Path(output_path), criteria)
    clean = remaining.total_matches == 0
    return VerificationResult(
        output_path=str(output_path), clean=clean, remaining=remaining if not clean else None
    )
